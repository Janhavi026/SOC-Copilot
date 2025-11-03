import os
import fitz # PyMuPDF for PDF
from docx import Document # python-docx for DOCX
from bs4 import BeautifulSoup # BeautifulSoup for HTML
from PIL import Image # Pillow for images
import pytesseract # Tesseract for OCR
import pandas as pd # pandas for CSV/Excel
import json
import logging
import yaml
import xml.etree.ElementTree as ET
from pptx import Presentation # python-pptx for PowerPoint
import markdown
from striprtf.striprtf import rtf_to_text # striprtf for RTF files

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def _load_text_from_pdf(file_path: str) -> str:
    """Extracts text from a PDF file."""
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        logging.info(f"Successfully extracted text from PDF: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from PDF {file_path}: {e}")
    return text

def _load_text_from_docx(file_path: str) -> str:
    """Extracts text from a DOCX file."""
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        logging.info(f"Successfully extracted text from DOCX: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from DOCX {file_path}: {e}")
    return text

def _load_text_from_html(file_path: str) -> str:
    """Extracts text from an HTML file."""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            text = soup.get_text(separator=' ', strip=True)
        logging.info(f"Successfully extracted text from HTML: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from HTML {file_path}: {e}")
    return text

def _load_text_from_txt(file_path: str) -> str:
    """Extracts text from a TXT file."""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        logging.info(f"Successfully extracted text from TXT: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from TXT {file_path}: {e}")
    return text

def _load_text_from_image(file_path: str) -> str:
    """Extracts text from an image file using OCR (Tesseract)."""
    text = ""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        logging.info(f"Successfully extracted text from image (OCR): {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from image {file_path} (OCR failed): {e}")
    return text

def _load_text_from_excel(file_path: str) -> str:
    """Extracts text from an Excel file (all sheets)."""
    text = ""
    try:
        xls = pd.ExcelFile(file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.to_string(index=False) + "\n"
        logging.info(f"Successfully extracted text from Excel: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from Excel {file_path}: {e}")
    return text

def _load_text_from_csv(file_path: str) -> str:
    """Extracts text from a CSV file."""
    text = ""
    try:
        df = pd.read_csv(file_path)
        text = df.to_string(index=False)
        logging.info(f"Successfully extracted text from CSV: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from CSV {file_path}: {e}")
    return text

def _format_product_info(product_data: dict) -> str:
    """
    Formats a single product's data into a readable string.
    This is specifically tailored for your 'products' array structure.
    """
    if not isinstance(product_data, dict):
        return ""

    product_str_parts = []
    product_str_parts.append(f"Product Name: {product_data.get('name', 'N/A')}")
    product_str_parts.append(f"Product ID: {product_data.get('product_id', 'N/A')}")
    product_str_parts.append(f"Type: {product_data.get('type', 'N/A')}")
    product_str_parts.append(f"Tier: {product_data.get('tier', 'N/A')}")

    price = product_data.get('price')
    price_unit = product_data.get('price_unit')
    if price is not None and price_unit:
        product_str_parts.append(f"Price: {price} {price_unit}")
    elif price is not None:
        product_str_parts.append(f"Price: {price}")
    
    if product_data.get('target_audience'):
        product_str_parts.append(f"Target Audience: {product_data['target_audience']}")
    if product_data.get('coverage'):
        product_str_parts.append(f"Coverage: {product_data['coverage']}")
    
    features = product_data.get('features', {})
    if features:
        product_str_parts.append("Features:")
        for category, detail_features in features.items():
            product_str_parts.append(f"  {category}:")
            for feature_name, has_feature in detail_features.items():
                if has_feature: # Only list features that are true
                    product_str_parts.append(f"    - {feature_name}")
    
    if product_data.get('compliance'):
        product_str_parts.append(f"Compliance: {', '.join(product_data['compliance'])}")
    if product_data.get('scalability'):
        product_str_parts.append(f"Scalability: {product_data['scalability']}")

    return "\n".join(product_str_parts)

def _load_text_from_json(file_path: str) -> str:
    """
    Extracts text from a JSON file, specifically handling
    the 'cybersecurity_products' structure for better semantic meaning.
    If it's a general JSON, it falls back to a flattening approach.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Check if it's your specific 'cybersecurity_products' format
        if isinstance(data, dict) and 'cybersecurity_products' in data and 'products' in data['cybersecurity_products']:
            logging.info(f"Detected 'cybersecurity_products' structure in {file_path}. Formatting product info.")
            all_product_text = []
            for product in data['cybersecurity_products']['products']:
                all_product_text.append(_format_product_info(product))
            
            # Also add feature descriptions if they exist
            if 'feature_descriptions' in data['cybersecurity_products']:
                all_product_text.append("\nFeature Descriptions:")
                for feature_name, description in data['cybersecurity_products']['feature_descriptions'].items():
                    all_product_text.append(f"- {feature_name}: {description}")

            return "\n\n".join(all_product_text)
        else:
            # Fallback for general JSONs: flatten keys and values
            logging.info(f"Detected general JSON structure in {file_path}. Flattening content.")
            def _extract_strings_from_general_json(obj):
                strings = []
                if isinstance(obj, dict):
                    for k, v in obj.items():
                        strings.append(f"{k}: {_extract_strings_from_general_json(v)}")
                elif isinstance(obj, list):
                    for item in obj:
                        strings.append(_extract_strings_from_general_json(item))
                else:
                    strings.append(str(obj))
                return " ".join(strings)
            return _extract_strings_from_general_json(data)
            
    except json.JSONDecodeError as e:
        logging.error(f"Error decoding JSON from {file_path}: {e}")
        return ""
    except Exception as e:
        logging.error(f"Error processing JSON file {file_path}: {e}")
        return ""


def _load_text_from_pptx(file_path: str) -> str:
    """Extracts text from a PowerPoint file."""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        logging.info(f"Successfully extracted text from PPTX: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from PPTX {file_path}: {e}")
    return text

def _load_text_from_rtf(file_path: str) -> str:
    """Extracts text from an RTF file."""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            rtf_content = f.read()
            text = rtf_to_text(rtf_content)
        logging.info(f"Successfully extracted text from RTF: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from RTF {file_path}: {e}")
    return text

def _load_text_from_xml(file_path: str) -> str:
    """Extracts text from an XML file."""
    text = ""
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        def extract_text_from_element(element):
            result = []
            if element.text:
                result.append(element.text.strip())
            for child in element:
                result.extend(extract_text_from_element(child))
            if element.tail:
                result.append(element.tail.strip())
            return result
        
        all_text = extract_text_from_element(root)
        text = " ".join([t for t in all_text if t])
        logging.info(f"Successfully extracted text from XML: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from XML {file_path}: {e}")
    return text

def _load_text_from_yaml(file_path: str) -> str:
    """Extracts text from a YAML file."""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = yaml.safe_load(f)
        
        def extract_yaml_content(obj, prefix=""):
            result = []
            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, (dict, list)):
                        result.append(f"{prefix}{key}:")
                        result.extend(extract_yaml_content(value, prefix + "  "))
                    else:
                        result.append(f"{prefix}{key}: {value}")
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    result.extend(extract_yaml_content(item, prefix + f"[{i}] "))
            else:
                result.append(f"{prefix}{obj}")
            return result
        
        text = "\n".join(extract_yaml_content(data))
        logging.info(f"Successfully extracted text from YAML: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from YAML {file_path}: {e}")
    return text

def _load_text_from_markdown(file_path: str) -> str:
    """Extracts text from a Markdown file."""
    text = ""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        # Convert markdown to HTML then extract text
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, 'html.parser')
        text = soup.get_text(separator=' ', strip=True)
        logging.info(f"Successfully extracted text from Markdown: {file_path}")
    except Exception as e:
        logging.error(f"Error extracting text from Markdown {file_path}: {e}")
    return text


def load_all_documents_from_directory(directory_path: str) -> list[dict]:
    """
    Loads text content from all supported document types in a directory
    and its subdirectories.
    Returns a list of dictionaries, each with 'content' and 'file_path'.
    """
    all_documents = []
    supported_extensions = {
        '.pdf': _load_text_from_pdf,
        '.docx': _load_text_from_docx,
        '.pptx': _load_text_from_pptx,
        '.ppt': _load_text_from_pptx,  # Note: requires additional handling for .ppt
        '.rtf': _load_text_from_rtf,
        '.html': _load_text_from_html,
        '.htm': _load_text_from_html,
        '.xml': _load_text_from_xml,
        '.yaml': _load_text_from_yaml,
        '.yml': _load_text_from_yaml,
        '.md': _load_text_from_markdown,
        '.markdown': _load_text_from_markdown,
        '.txt': _load_text_from_txt,
        '.png': _load_text_from_image,
        '.jpg': _load_text_from_image,
        '.jpeg': _load_text_from_image,
        '.webp': _load_text_from_image,
        '.json': _load_text_from_json,
        '.csv': _load_text_from_csv,
        '.xlsx': _load_text_from_excel,
    }

    if not os.path.exists(directory_path):
        logging.warning(f"Directory not found: {directory_path}")
        return []

    for root, _, files in os.walk(directory_path):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            file_extension = os.path.splitext(file_name)[1].lower()

            if file_extension in supported_extensions:
                logging.info(f"Attempting to load: {file_path}")
                content = supported_extensions[file_extension](file_path)
                if content:
                    all_documents.append({"content": content, "file_path": file_path})
                else:
                    logging.warning(f"No content extracted or error during loading for {file_path}")
            else:
                logging.info(f"Skipping unsupported file type: {file_path}")
    
    if not all_documents:
        logging.warning(f"No documents loaded from {directory_path}. Check if directory contains supported files.")

    return all_documents